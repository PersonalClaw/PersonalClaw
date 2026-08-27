"""DFE-8 — the deck round trip: a bullet keeps its depth.

**Every assertion here re-parses the WRITTEN BYTES.** A test that renders a model and
compares it to itself proves only that the writer is self-consistent; the claim under test
is that the *file* carries the depth, so each case goes model → bytes →
:func:`parse_pptx` → model, and several go one step further and read the bytes with
python-pptx directly so the evidence does not depend on our own parser either.

The defect this atom closes: the writer pinned ``para.level = 0`` on every appended
paragraph, so a nested outline came out flat. It was not lost in the file — it was never
written, and nothing ever read one back, because no .pptx parser existed at all
(``GET …/model`` answered 415 for every deck).
"""

from __future__ import annotations

import io

import pytest

from personalclaw.documents.deck_json import deck_from_dict, deck_to_dict
from personalclaw.documents.model import (
    DECK_LAYOUTS,
    MAX_BULLET_LEVEL,
    Bullet,
    DeckModel,
    ShapeBox,
    Slide,
)
from personalclaw.documents.model_codec import get_codec
from personalclaw.documents.pptx_parser import parse_pptx
from personalclaw.documents.writers.pptx_writer import render_pptx


def _lap(model: DeckModel) -> DeckModel:
    """One full trip through the real writer AND the real parser."""
    parsed, _ = parse_pptx(render_pptx(model))
    return parsed


def _losses(model: DeckModel):
    return parse_pptx(render_pptx(model))[1]


def _slides(data: bytes):
    """The written file's slides, read with python-pptx — not with our parser."""
    from pptx import Presentation

    return list(Presentation(io.BytesIO(data)).slides)


def _outline(model: DeckModel, index: int = 0) -> list[tuple[str, int]]:
    return [(b.text, b.level) for b in model.slides[index].bullets]


# ── the headline: bullet depth ────────────────────────────────────────────────


def test_bullet_depth_survives_the_round_trip() -> None:
    """The atom's first clause. Depth in, depth out, through the real writer and parser."""
    deck = DeckModel(
        slides=[
            Slide(
                title="Plan",
                bullets=[
                    Bullet(text="top"),
                    Bullet(text="under it", level=1),
                    Bullet(text="deeper still", level=2),
                    Bullet(text="back up", level=1),
                ],
            )
        ]
    )
    assert _outline(_lap(deck)) == [
        ("top", 0),
        ("under it", 1),
        ("deeper still", 2),
        ("back up", 1),
    ]


def test_the_written_file_itself_carries_the_depth() -> None:
    """The same claim without our parser in the loop: python-pptx reads ``a:pPr/@lvl``
    straight out of the bytes, so this is the file's own account of the outline."""
    data = render_pptx(
        DeckModel(
            slides=[
                Slide(
                    title="Plan",
                    bullets=[Bullet(text="a"), Bullet(text="b", level=3), Bullet(text="c")],
                )
            ]
        )
    )
    from personalclaw.documents.pptx_shapes import body_placeholder

    body = body_placeholder(_slides(data)[0])
    assert body is not None
    assert [(p.text, p.level) for p in body.text_frame.paragraphs] == [
        ("a", 0),
        ("b", 3),
        ("c", 0),
    ]


def test_the_first_bullet_keeps_its_depth_too() -> None:
    """The writer reuses the placeholder's EXISTING first paragraph rather than appending
    one, so the first bullet takes a different code path from the rest — and a fix applied
    only to the appended ones would leave a whole deck's opening line flat."""
    deck = DeckModel(slides=[Slide(title="T", bullets=[Bullet(text="indented", level=2)])])
    assert _outline(_lap(deck)) == [("indented", 2)]


def test_a_flat_outline_stays_flat() -> None:
    """Vacuity for the depth tests: a writer that wrote a fixed non-zero level, or a parser
    that invented depth from position, would fail here while passing the ones above."""
    deck = DeckModel(slides=[Slide.outline("T", ["one", "two", "three"])])
    assert _outline(_lap(deck)) == [("one", 0), ("two", 0), ("three", 0)]


def test_a_depth_beyond_what_powerpoint_can_express_is_clamped() -> None:
    """The format's own schema stops at 8. Clamped rather than refused: the caller has an
    outline, and the deepest indent is what a reader would have seen anyway."""
    assert Bullet(text="x", level=40).level == MAX_BULLET_LEVEL
    assert Bullet(text="x", level=-3).level == 0
    assert _outline(_lap(DeckModel(slides=[Slide(bullets=[Bullet(text="x", level=40)])]))) == [
        ("x", MAX_BULLET_LEVEL)
    ]


def test_a_second_lap_changes_nothing() -> None:
    """Idempotence is the property an editor needs: opening a deck and saving it without
    touching anything must not drift the file. A parse that re-derived depth, geometry or
    the layout from a default would show up here as the second lap disagreeing."""
    deck = DeckModel(
        title="Cover",
        slides=[
            Slide(title="A", bullets=[Bullet(text="a"), Bullet(text="b", level=1)], notes="n"),
            Slide(title="B", layout="Section Header"),
        ],
    )
    once = _lap(deck)
    assert _lap(once) == once


# ── layout and geometry ──────────────────────────────────────────────────────


def test_the_declared_layouts_are_the_shipped_templates_own() -> None:
    """``DECK_LAYOUTS`` is the editor's option list and the parser's "can we re-create
    this?" test, so a template whose layouts were renamed must fail HERE rather than
    silently re-laying out every slide it names."""
    from pptx import Presentation

    assert tuple(layout.name for layout in Presentation().slide_layouts) == DECK_LAYOUTS


def test_a_named_layout_round_trips_and_reaches_the_file() -> None:
    deck = DeckModel(slides=[Slide(title="Part two", layout="Section Header")])
    assert _lap(deck).slides[0].layout == "Section Header"
    assert _slides(render_pptx(deck))[0].slide_layout.name == "Section Header"


def test_a_slide_with_no_named_layout_is_laid_out_from_its_content() -> None:
    """Vacuity for the test above — and the behaviour a generated deck relies on: an
    empty body placeholder renders as a "Click to add text" prompt, so a slide with no
    bullets must not get the content layout."""
    deck = DeckModel(slides=[Slide(title="Bare"), Slide.outline("Full", ["x"])])
    assert [s.layout for s in _lap(deck).slides] == ["Title Only", "Title and Content"]


def test_a_moved_shape_keeps_its_position_through_the_round_trip() -> None:
    box = ShapeBox(left_in=1.25, top_in=0.5, width_in=6.0, height_in=1.5)
    deck = DeckModel(slides=[Slide(title="Moved", bullets=[Bullet(text="b")], title_box=box)])
    got = _lap(deck).slides[0].title_box
    assert got.placed
    assert (got.left_in, got.top_in, got.width_in, got.height_in) == pytest.approx(
        (1.25, 0.5, 6.0, 1.5), abs=1e-3
    )


def test_a_shape_that_never_moved_stays_INHERITED_rather_than_pinned() -> None:
    """The discriminating half of the geometry story. python-pptx reports a placeholder's
    inherited position as if the shape declared it, so a parser that stored what it read
    would pin every shape of every deck at whatever the template said that day — and the
    next template change would stop reaching the file."""
    parsed = _lap(DeckModel(slides=[Slide(title="Still", bullets=[Bullet(text="b")])]))
    assert not parsed.slides[0].title_box.placed
    assert not parsed.slides[0].body_box.placed
    assert parsed.slides[0].title_box == ShapeBox()


def test_a_widescreen_slide_size_round_trips() -> None:
    deck = DeckModel(width_in=13.333, height_in=7.5, slides=[Slide(title="Wide")])
    parsed = _lap(deck)
    assert (parsed.width_in, parsed.height_in) == pytest.approx((13.333, 7.5), abs=1e-3)


def test_the_templates_own_slide_size_is_what_an_unset_size_means() -> None:
    """Vacuity for the size test: ``0.0`` is "the template's size", not "zero"."""
    assert (_lap(DeckModel(slides=[Slide(title="T")])).width_in) == 10.0


# ── the cover slide ──────────────────────────────────────────────────────────


def test_a_deck_title_round_trips_through_its_cover_slide() -> None:
    """The writer renders a deck title AS a title slide, so the parser reads that shape
    back into the title field. Without this, saving an edited title would append a SECOND
    cover on every save — the field and the slide list would both grow."""
    parsed = _lap(DeckModel(title="Quarterly Review", slides=[Slide.outline("A", ["x"])]))
    assert parsed.title == "Quarterly Review"
    assert [s.title for s in parsed.slides] == ["A"]


def test_a_first_slide_that_carries_content_is_NOT_folded_into_the_title() -> None:
    """Folding removes a slide from the list, so it is only safe when nothing on it would
    have nowhere to go. A title slide with bullets stays a slide."""
    deck = DeckModel(
        slides=[Slide(title="Cover-ish", layout="Title Slide", bullets=[Bullet(text="agenda")])]
    )
    parsed = _lap(deck)
    assert parsed.title == ""
    assert [(s.title, s.layout) for s in parsed.slides] == [("Cover-ish", "Title Slide")]


def test_notes_and_the_image_reference_survive_the_lap() -> None:
    deck = DeckModel(slides=[Slide(title="Chart", notes="mention Q3", artifact_slug="sales")])
    notes = _lap(deck).slides[0].notes
    assert "mention Q3" in notes
    # The writer records an unresolvable image reference in the notes rather than dropping
    # it. The parser does NOT sniff it back out into `artifact_slug`: recovering a field by
    # matching a text prefix is the guessing this family exists to abolish.
    assert "[image: sales]" in notes


def test_a_trailing_empty_paragraph_is_not_read_back_as_a_bullet() -> None:
    """PowerPoint leaves an empty final paragraph behind whenever a placeholder is touched
    and emptied; carrying it would grow a stray bullet glyph on every save."""
    from pptx import Presentation

    from personalclaw.documents.pptx_shapes import body_placeholder

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "T"
    frame = body_placeholder(slide).text_frame
    frame.text = "real"
    frame.add_paragraph()  # the trailing empty one
    buf = io.BytesIO()
    prs.save(buf)

    model, report = parse_pptx(buf.getvalue())
    assert _outline(model) == [("real", 0)]
    assert report.lossless, report.summary()


# ── the loss report ──────────────────────────────────────────────────────────


def test_a_plain_deck_is_lossless() -> None:
    """The shared vacuity leg for every loss test below: a deck this model can hold fully
    must report NOTHING, or a report of one item would mean nothing."""
    deck = DeckModel(
        title="Cover",
        slides=[
            Slide(title="A", bullets=[Bullet(text="a"), Bullet(text="b", level=1)], notes="n"),
            Slide(title="B"),
        ],
    )
    report = _losses(deck)
    assert report.lossless, report.summary()


def _deck_with(layout: int = 1, title: str = "T"):
    """A presentation with one slide on *layout*, for the loss cases to disfigure."""
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[layout])
    if slide.shapes.title is not None:
        slide.shapes.title.text = title
    return prs, slide


def _bytes(prs) -> bytes:
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_a_layout_the_shipped_template_cannot_re_create_is_reported() -> None:
    prs, slide = _deck_with()
    # A deck built from a corporate template carries its own layout names. Renamed through
    # the XML because python-pptx exposes the name read-only — this is a FIXTURE standing
    # in for a real template, not a mechanism under test.
    slide.slide_layout._element.cSld.set("name", "Acme Section Break")
    model, report = parse_pptx(_bytes(prs))

    assert [item.kind for item in report.items] == ["slide_layout"]
    assert report.items[0].where == "slide 1"
    assert "Acme Section Break" in report.items[0].detail
    # …and the slide is still editable, on no declared layout, rather than being refused.
    assert model.slides[0].layout == ""
    assert model.slides[0].title == "T"


def test_a_placeholder_the_model_has_no_field_for_is_reported_with_its_text() -> None:
    """A two-content layout has a second body. The model holds ONE outline, so the second
    one's text is reported — and quoted, so it is findable rather than merely counted."""
    prs, slide = _deck_with(layout=3)
    bodies = [
        ph for ph in slide.placeholders if ph.placeholder_format.idx != 0 and ph.has_text_frame
    ]
    bodies[0].text_frame.text = "carried"
    bodies[1].text_frame.text = "the right hand column"
    model, report = parse_pptx(_bytes(prs))

    assert [item.kind for item in report.items] == ["slide_placeholder"]
    assert "the right hand column" in report.items[0].detail
    assert report.items[0].where == "slide 1"
    assert _outline(model) == [("carried", 0)]


def test_an_unused_placeholder_is_not_reported() -> None:
    """Vacuity for the test above. Every layout ships empty placeholders (a Title Slide has
    a subtitle box); a report that named each one would be noise a user learns to ignore,
    which is how the item that matters gets missed."""
    prs, _slide = _deck_with(layout=3)
    assert parse_pptx(_bytes(prs))[1].lossless


def test_a_free_shape_is_reported_not_dropped_silently() -> None:
    from pptx.util import Inches

    prs, slide = _deck_with()
    box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(3), Inches(1))
    box.text_frame.text = "a floating caption"
    report = parse_pptx(_bytes(prs))[1]

    assert [item.kind for item in report.items] == ["slide_shape"]
    assert "a floating caption" in report.items[0].detail
    assert report.items[0].where == "slide 1"


def test_character_formatting_inside_a_bullet_is_reported_at_that_bullet() -> None:
    """Located per BULLET, not per slide: "somewhere on slide 4 there was a bold word" is
    not something a person can act on."""
    from personalclaw.documents.pptx_shapes import body_placeholder

    prs, slide = _deck_with()
    frame = body_placeholder(slide).text_frame
    frame.text = "plain"
    second = frame.add_paragraph()
    second.text = "shouty"
    second.runs[0].font.bold = True
    third = frame.add_paragraph()
    third.text = "linked"
    third.runs[0].hyperlink.address = "https://example.invalid/x"
    report = parse_pptx(_bytes(prs))[1]

    assert [item.kind for item in report.items] == ["bullet_run_style", "bullet_run_style"]
    assert [item.where for item in report.items] == ["slide 1 · bullet 2", "slide 1 · bullet 3"]
    assert "bold" in report.items[0].detail
    assert "https://example.invalid/x" in report.items[1].detail


def test_a_bullet_whose_formatting_is_inherited_is_not_reported() -> None:
    """Vacuity: ``None`` means "inherited from the layout", which is the normal case. A
    parser that read None as a style would put an item on every bullet of every deck."""
    prs, slide = _deck_with()
    from personalclaw.documents.pptx_shapes import body_placeholder

    body_placeholder(slide).text_frame.text = "plain"
    assert parse_pptx(_bytes(prs))[1].lossless


def test_a_slide_that_overrides_its_background_is_reported() -> None:
    from pptx.dml.color import RGBColor

    prs, slide = _deck_with()
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0x22, 0x33, 0x44)
    report = parse_pptx(_bytes(prs))[1]

    assert [item.kind for item in report.items] == ["slide_feature"]
    assert report.items[0].where == "slide 1"


def test_a_loss_on_a_folded_cover_is_located_as_the_cover() -> None:
    """A loss is located in the terms the EDITOR shows, which is the whole point of
    `LossItem.location`: the cover is a field, not slide 1, so calling it "slide 1" would
    send a user looking at the wrong surface."""
    from pptx.util import Inches

    prs, slide = _deck_with(layout=0, title="Cover")
    slide.shapes.add_textbox(Inches(1), Inches(6), Inches(2), Inches(0.5))
    model, report = parse_pptx(_bytes(prs))

    assert model.title == "Cover"
    assert [(item.kind, item.where) for item in report.items] == [("slide_shape", "cover slide")]


def test_a_title_slides_strapline_is_carried_as_its_body_not_lost() -> None:
    """MEASURED, and it changes what a cover means: a Title Slide's subtitle IS the first
    text-frame placeholder that is not the title, so the model carries it as the slide's
    outline and the writer puts it back there. A strapline is therefore not a loss — and a
    cover carrying one keeps its own slide rather than folding into the title field, which
    is the same rule as bullets: folding is only safe when nothing is left over."""
    prs, slide = _deck_with(layout=0, title="Cover")
    subtitle = [ph for ph in slide.placeholders if ph.placeholder_format.idx == 1][0]
    subtitle.text_frame.text = "a strapline"
    model, report = parse_pptx(_bytes(prs))

    assert report.lossless, report.summary()
    assert model.title == ""
    assert _outline(model) == [("a strapline", 0)]
    assert _lap(model).slides[0].bullets[0].text == "a strapline"


# ── markdown → depth ─────────────────────────────────────────────────────────


def test_markdown_indentation_becomes_bullet_depth() -> None:
    """Depth read from markdown's OWN encoding of it — leading whitespace — which
    `_BULLET`'s `^\\s*` used to consume and throw away."""
    from personalclaw.documents.from_markup import deck_from_markdown

    deck = deck_from_markdown("## Plan\n- top\n  - under\n    - deeper\n- back\n")
    assert [(b.text, b.level) for b in deck.slides[0].bullets] == [
        ("top", 0),
        ("under", 1),
        ("deeper", 2),
        ("back", 0),
    ]
    # …and it reaches the FILE, not just the model.
    assert _outline(_lap(deck)) == [("top", 0), ("under", 1), ("deeper", 2), ("back", 0)]


def test_an_unindented_outline_is_all_top_level() -> None:
    from personalclaw.documents.from_markup import deck_from_markdown

    deck = deck_from_markdown("## Plan\n- one\n- two\n")
    assert [b.level for b in deck.slides[0].bullets] == [0, 0]


# ── the wire boundary ────────────────────────────────────────────────────────


def test_a_deck_survives_the_json_boundary_unchanged() -> None:
    deck = DeckModel(
        title="C",
        width_in=13.333,
        height_in=7.5,
        slides=[
            Slide(
                title="A",
                bullets=[Bullet(text="a", level=1)],
                notes="n",
                layout="Title and Content",
                body_box=ShapeBox(left_in=1.0, top_in=2.0, width_in=3.0, height_in=4.0),
            )
        ],
    )
    assert deck_from_dict(deck_to_dict(deck)) == deck


def test_an_unknown_field_is_refused_rather_than_dropped() -> None:
    payload = deck_to_dict(DeckModel(slides=[Slide(title="A")]))
    payload["slides"][0]["bulletts"] = []
    with pytest.raises(ValueError, match="unknown field"):
        deck_from_dict(payload)


def test_a_layout_the_template_lacks_is_refused_at_the_wire() -> None:
    """Refused, not dropped: a client that sent ``"Titel Slide"`` would otherwise get a
    deck silently re-laid-out from its content and find out by opening the file."""
    payload = deck_to_dict(DeckModel(slides=[Slide(title="A")]))
    payload["slides"][0]["layout"] = "Titel Slide"
    with pytest.raises(ValueError, match="model.slides\\[0\\].layout"):
        deck_from_dict(payload)


def test_a_bullet_depth_that_is_not_a_number_is_refused() -> None:
    payload = deck_to_dict(DeckModel(slides=[Slide(bullets=[Bullet(text="a")])]))
    payload["slides"][0]["bullets"][0]["level"] = "2"
    with pytest.raises(ValueError, match="must be an integer"):
        deck_from_dict(payload)


def test_a_depth_out_of_range_from_the_wire_is_clamped_not_refused() -> None:
    payload = deck_to_dict(DeckModel(slides=[Slide(bullets=[Bullet(text="a")])]))
    payload["slides"][0]["bullets"][0]["level"] = 99
    assert deck_from_dict(payload).slides[0].bullets[0].level == MAX_BULLET_LEVEL


def test_a_server_side_caller_cannot_name_a_layout_the_writer_cannot_resolve() -> None:
    """The wire's guard has a twin on the dataclass, so a Python caller assembling a deck
    gets the same refusal rather than a slide that quietly rebuilt itself."""
    with pytest.raises(ValueError, match="unknown layout"):
        Slide(title="A", layout="Wonky")


# ── the codec table ──────────────────────────────────────────────────────────


def test_the_pptx_codec_is_wired_to_the_shipped_parser_and_writer() -> None:
    """The call site, not the mechanism: the codec the ROUTE looks up must be the parser
    and serializer this suite proved, not a second pair."""
    codec = get_codec("pptx")
    assert codec is not None
    assert codec.parse is parse_pptx
    assert codec.to_dict is deck_to_dict
    assert codec.from_dict is deck_from_dict
