"""File readers for knowledge ingestion. Supports text, PDF, PPTX, DOCX, HTML, XLSX."""

import os
import re
from dataclasses import dataclass, field
from pathlib import Path

from personalclaw.security import is_sensitive_path

try:
    import pdfplumber
except ImportError:
    pdfplumber = None  # type: ignore[assignment]

try:
    from pptx import Presentation  # type: ignore[import-untyped]
except ImportError:
    Presentation = None  # type: ignore[assignment,misc]

try:
    from docx import Document  # type: ignore[import-untyped]
except ImportError:
    Document = None  # type: ignore[assignment,misc]

try:
    import html2text as _html2text_mod
except ImportError:
    _html2text_mod = None  # type: ignore[assignment]

try:
    from openpyxl import load_workbook as _load_workbook  # type: ignore[import-untyped]
except ImportError:
    _load_workbook = None  # type: ignore[assignment]


def _render_docx_table(table) -> list[str]:
    """One .docx table → markdown rows. Empty list when the table holds nothing.

    Merged cells repeat their text in python-docx (the same cell object appears at each
    covered position); de-duplicating adjacent repeats keeps a merged header from reading
    as "Q1 | Q1 | Q1" without inventing a colspan notation markdown cannot express.
    """
    rows: list[list[str]] = []
    for row in table.rows:
        cells: list[str] = []
        for cell in row.cells:
            text = " ".join(cell.text.split())
            # Pipes would break the markdown table structure.
            text = text.replace("|", "\\|")
            if cells and text and text == cells[-1]:
                continue  # a merged cell repeating itself
            cells.append(text)
        if any(c for c in cells):
            rows.append(cells)
    if not rows:
        return []
    width = max(len(r) for r in rows)
    out = ["| " + " | ".join((rows[0] + [""] * width)[:width]) + " |"]
    out.append("| " + " | ".join(["---"] * width) + " |")
    for row in rows[1:]:
        out.append("| " + " | ".join((row + [""] * width)[:width]) + " |")
    return out


class FileReader:
    # Note: .pdf/.pptx require optional runtime dependencies (pdfplumber, python-pptx)
    # not declared in setup.cfg. .docx requires python-docx (pip install python-docx).
    SUPPORTED = {
        "",
        ".md",
        ".markdown",
        ".txt",
        ".text",
        ".py",
        ".java",
        ".ts",
        ".js",
        ".rs",
        ".go",
        ".html",
        ".htm",
        ".docx",
        ".csv",
        ".tsv",
        ".log",
        ".json",
        ".yaml",
        ".yml",
        ".sh",
        ".rb",
        ".c",
        ".cpp",
        ".h",
        ".xlsx",
        ".xls",
    }

    _DISPATCH = {
        ".pdf": "_read_pdf",
        ".pptx": "_read_pptx",
        ".docx": "_read_docx",
        ".html": "_read_html",
        ".htm": "_read_html",
        ".xlsx": "_read_xlsx",
        ".xls": "_read_xlsx",
        ".csv": "_read_csv",
        ".tsv": "_read_csv",
    }

    # Cap the markdown-table rendering of a large CSV so a huge file doesn't bloat the
    # stored content / embedding; the true row_count is still recorded in metadata.
    _CSV_MAX_TABLE_ROWS = 500

    def read(self, path: str) -> tuple[str, dict]:
        if is_sensitive_path(path):
            raise PermissionError(f"Refusing to read sensitive path: {path}")
        p = Path(path)
        ext = p.suffix.lower()
        base_meta = {
            "format": ext.lstrip("."),
            "title": p.stem,
            "file_size": os.path.getsize(path),
            "extension": ext,
        }
        method_name = self._DISPATCH.get(ext)
        if method_name:
            text, meta = getattr(self, method_name)(path)
            base_meta.update(meta)
        else:
            text, meta = self._read_text(path, ext.lstrip("."))
            base_meta.update(meta)
        base_meta["line_count"] = text.count("\n") + 1 if text else 0
        return text, base_meta

    def _read_text(self, path: str, fmt: str) -> tuple[str, dict]:
        try:
            try:
                with open(path, "r", encoding="utf-8") as f:
                    text = f.read()
            except UnicodeDecodeError:
                with open(path, "r", encoding="latin-1") as f:
                    text = f.read()
            return text, {"format": fmt}
        except Exception as e:
            return f"Error reading file: {e}", {"format": "error", "error": str(e)}

    def _read_pdf(self, path: str) -> tuple[str, dict]:
        if pdfplumber is None:
            return (
                "PDF support requires pdfplumber: pip install pdfplumber",
                {"format": "error", "error": "PDF support requires pdfplumber"},
            )
        try:
            with pdfplumber.open(path) as pdf:
                pages = [p.extract_text() or "" for p in pdf.pages]
                return "\n".join(pages), {"format": "pdf", "page_count": len(pages)}
        except Exception as e:
            # A .pdf that isn't a real PDF is often a mislabeled text/markdown file.
            # Salvage it: read as text when the bytes decode to mostly-printable
            # content, rather than failing outright and losing the user's content.
            salvaged = self._salvage_as_text(path)
            if salvaged is not None:
                return salvaged, {"format": "text", "recovered_from": "pdf"}
            return f"Error reading file: {e}", {"format": "error", "error": str(e)}

    @staticmethod
    def _salvage_as_text(path: str) -> str | None:
        """Return decoded text if the file is plausibly plain text (mostly printable),
        else None. Used to recover a mislabeled text file from a failed binary parse."""
        try:
            with open(path, "rb") as f:
                raw = f.read(200_000)
        except OSError:
            return None
        if not raw:
            return None
        try:
            text = raw.decode("utf-8")
        except UnicodeDecodeError:
            return None
        sample = text[:2000]
        printable = sum(1 for c in sample if c.isprintable() or c in "\n\r\t ")
        # Require the sample to be ≥90% printable to treat it as genuine text.
        if printable / max(1, len(sample)) < 0.9:
            return None
        return text

    def _read_pptx(self, path: str) -> tuple[str, dict]:
        if Presentation is None:
            return (
                "PPTX support requires python-pptx: pip install python-pptx",
                {"format": "error", "error": "PPTX support requires python-pptx"},
            )
        try:
            prs = Presentation(path)
            parts = []
            for i, slide in enumerate(prs.slides, 1):
                title = ""
                body_parts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text_frame.text.strip()
                        if shape == slide.shapes.title:
                            title = text
                        else:
                            body_parts.append(text)
                notes = ""
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                section = f"## Slide {i}: {title}\n{chr(10).join(body_parts)}"
                if notes:
                    section += f"\n{notes}"
                parts.append(section)
            return "\n\n".join(parts), {"format": "pptx", "slide_count": len(prs.slides)}
        except Exception as e:
            return f"Error reading file: {e}", {"format": "error", "error": str(e)}

    def _read_docx(self, path: str) -> tuple[str, dict]:
        if Document is None:
            return (
                "DOCX support requires python-docx: pip install python-docx",
                {"format": "error", "error": "DOCX support requires python-docx"},
            )
        try:
            doc = Document(path)
            lines = []
            for para in doc.paragraphs:
                style = para.style.name if para.style else ""
                text = para.text
                if style.startswith("Heading"):
                    try:
                        level = int(style.split()[-1])
                    except (ValueError, IndexError):
                        level = 1
                    lines.append(f'{"#" * level} {text}')
                else:
                    lines.append(text)
            # Tables were previously DROPPED ENTIRELY: this only walked
            # `doc.paragraphs`, and python-docx keeps table content out of that
            # collection. Any table in an ingested Word document — often the densest
            # information in it — was silently invisible to search, embedding and the
            # agent. Rendered as markdown tables to match how the xlsx/csv readers
            # already present tabular data, so downstream consumers see one shape.
            #
            # Appended after the prose rather than interleaved, and that is now a CHOICE
            # rather than a limit: `documents/docx_parser._walk` walks the body element's
            # own children, which IS document order, so the ordering this reader does not
            # reproduce is available in this repo. It still appends because this path
            # flattens a document to plain text for search and embedding, where no
            # consumer reads position — and because interleaving here changes the stored
            # text of every already-ingested Word document. Losing POSITION was always a
            # far smaller defect than losing the content; converging on the parser's walk
            # belongs to whichever change is willing to own that re-ingest.
            for table in doc.tables:
                rendered = _render_docx_table(table)
                if rendered:
                    lines.append("")
                    lines.extend(rendered)
            return "\n".join(lines), {
                "format": "docx",
                "content_type": "markdown",
                "paragraph_count": len(doc.paragraphs),
                "table_count": len(doc.tables),
            }
        except Exception as e:
            return f"Error reading file: {e}", {"format": "error", "error": str(e)}

    def _read_xlsx(self, path: str) -> tuple[str, dict]:
        """Extract a spreadsheet as markdown tables (one per sheet). Without openpyxl,
        return an error rather than letting the binary .xlsx be read as raw text."""
        if _load_workbook is None:
            return (
                "XLSX support requires openpyxl: pip install openpyxl",
                {"format": "error", "error": "XLSX support requires openpyxl"},
            )
        try:
            wb = _load_workbook(path, read_only=True, data_only=True)
            parts, total_rows = [], 0
            for ws in wb.worksheets:
                rows = [
                    ["" if c is None else str(c) for c in row]
                    for row in ws.iter_rows(values_only=True)
                ]
                rows = [r for r in rows if any(cell.strip() for cell in r)]
                if not rows:
                    continue
                total_rows += len(rows)
                lines = [
                    f"## {ws.title}",
                    "| " + " | ".join(rows[0]) + " |",
                    "| " + " | ".join("---" for _ in rows[0]) + " |",
                ]
                lines += ["| " + " | ".join(r) + " |" for r in rows[1:]]
                parts.append("\n".join(lines))
            sheet_count = len(wb.worksheets)
            wb.close()
            return (
                "\n\n".join(parts),
                {
                    "format": "xlsx",
                    "content_type": "markdown",
                    "sheet_count": sheet_count,
                    "row_count": total_rows,
                },
            )
        except Exception as e:
            return f"Error reading spreadsheet: {e}", {"format": "error", "error": str(e)}

    def _read_csv(self, path: str) -> tuple[str, dict]:
        """Render a CSV/TSV as a markdown table (consistent with _read_xlsx), so a tabular
        upload — a 'sheet'-type item — ingests as structured content + row_count metadata
        rather than raw delimited text. Uses the csv module so quoted fields/embedded
        delimiters parse correctly. Large files render a capped table; row_count is true."""
        import csv as _csv

        # .tsv is tab-delimited; everything else (.csv) is comma-delimited.
        delimiter = "\t" if Path(path).suffix.lower() == ".tsv" else ","
        fmt = "tsv" if delimiter == "\t" else "csv"
        try:
            try:
                f = open(path, newline="", encoding="utf-8")
            except UnicodeDecodeError:
                f = open(path, newline="", encoding="latin-1")
            with f:
                rows = [
                    [("" if c is None else str(c)) for c in row]
                    for row in _csv.reader(f, delimiter=delimiter)
                ]
        except Exception as e:
            return f"Error reading {fmt.upper()}: {e}", {"format": "error", "error": str(e)}
        rows = [r for r in rows if any(cell.strip() for cell in r)]
        if not rows:
            return "", {"format": fmt, "content_type": "markdown", "row_count": 0}
        width = max(len(r) for r in rows)
        rows = [r + [""] * (width - len(r)) for r in rows]  # pad ragged rows
        shown = rows[: self._CSV_MAX_TABLE_ROWS]

        # Escape pipes so cell content can't break the markdown table layout.
        def cells(r: list[str]) -> str:
            return "| " + " | ".join(c.replace("|", "\\|") for c in r) + " |"

        lines = [cells(shown[0]), "| " + " | ".join("---" for _ in shown[0]) + " |"]
        lines += [cells(r) for r in shown[1:]]
        if len(rows) > len(shown):
            lines.append(f"\n_…{len(rows) - len(shown)} more rows_")
        return "\n".join(lines), {"format": fmt, "content_type": "markdown", "row_count": len(rows)}

    def _read_html(self, path: str) -> tuple[str, dict]:
        try:
            with open(path, "r", encoding="utf-8") as f:
                html = f.read()
        except UnicodeDecodeError:
            with open(path, "r", encoding="latin-1") as f:
                html = f.read()
        except Exception as e:
            return f"Error reading file: {e}", {"format": "error", "error": str(e)}
        return html_to_prose(html), {"format": "html"}


def html_to_prose(html: str) -> str:
    """HTML markup → the prose a reader would see, with chrome stripped.

    The conversion half of :meth:`FileReader._read_html`, lifted out because an artifact's
    body arrives as a STRING from the artifact store rather than as a file on disk (its
    on-disk name is always ``current.html`` whatever the artifact's kind, so the extension
    dispatch cannot be reused as-is). One shared primitive rather than two: an ``html``
    artifact and an uploaded ``.html`` must reduce to the same text, and a second
    implementation is where "the chrome is stripped for uploads but not for artifacts"
    comes from.

    Chrome (nav/header/footer/aside/script/…) is dropped BEFORE conversion — the same
    primitive as the bookmark scrape — and the regex path is the degradation when
    ``html2text`` is not installed, not a second policy.
    """
    from personalclaw.knowledge.connectors.base import strip_html_chrome

    html = strip_html_chrome(html or "")
    if _html2text_mod is not None:
        h = _html2text_mod.HTML2Text()
        h.ignore_links = False
        h.ignore_images = True
        return h.handle(html)
    text = re.sub(r"<[^>]+>", " ", html)
    return re.sub(r"\s+", " ", text).strip()


# ── Structural PDF read (WATCHED-SOURCES §5) ──────────────────────────────────
#
# `_read_pdf` above answers "what does this document SAY"; a section detector needs
# "how is it LAID OUT" — per-page text, per-line font size, and the PDF's own outline.
# Both answers come from the SAME `pdfplumber` import at the top of this module, on
# purpose: a second PDF library (or a second guarded import) would be a second place
# for the PDF path to be present-or-absent, and callers would then have to know which
# one degraded. `knowledge/slicing.py` consumes these types and never imports
# pdfplumber itself, which is also what keeps the slicing cascade a pure function of
# a plain dataclass — testable with no PDF at all.


@dataclass(frozen=True)
class PdfLine:
    """One visual line of a PDF page: its text and the LARGEST glyph size on it.

    Max rather than mean because a heading followed by a superscript footnote marker
    must still read as heading-sized; averaging would drag it toward body size.
    """

    page: int  # 0-based page index
    text: str
    size: float
    char_count: int


@dataclass(frozen=True)
class PdfStructure:
    """A PDF's layout, as much as a text extractor can see.

    ``outline`` holds the document's own bookmark TITLES in document order and NOT
    their destinations: resolving a PDF destination to a page index goes through
    named-destination indirection that many real papers get wrong, whereas a title
    can simply be LOCATED in the extracted text — which is the offset the caller
    needs anyway. ``outline`` is empty for the common case of a paper with no
    bookmarks at all (measured: reportlab-generated PDFs carry none).
    """

    pages: tuple[str, ...] = ()
    lines: tuple[PdfLine, ...] = ()
    outline: tuple[str, ...] = field(default=())


def read_pdf_structure(path: str) -> PdfStructure | None:
    """Extract *path*'s layout, or None when it cannot be read as a PDF.

    None (rather than an empty structure) is deliberate: "this is not a PDF we can
    lay out" and "this is a PDF with one blank page" are different facts, and a caller
    that must fall back to text-only detection needs to tell them apart.
    """
    if pdfplumber is None:
        return None
    if is_sensitive_path(path):
        raise PermissionError(f"Refusing to read sensitive path: {path}")
    try:
        with pdfplumber.open(path) as pdf:
            pages = tuple((p.extract_text() or "") for p in pdf.pages)
            lines: list[PdfLine] = []
            for index, page in enumerate(pdf.pages):
                lines.extend(_lines_for_page(page, index))
            return PdfStructure(pages=pages, lines=tuple(lines), outline=_outline_titles(pdf))
    except Exception:
        # Mirrors `_read_pdf`'s tolerance: a mislabeled or damaged file must not raise
        # into an ingest. The caller degrades to text-only section detection.
        return None


#: Two glyph baselines within this many points are treated as the same visual line.
#: PDF writers jitter the `top` of glyphs on one line by fractions of a point (and
#: rounding to whole points would merge two lines of a tightly-leaded paper), so the
#: grouping needs a tolerance rather than an equality test.
#:
#: Deliberately NOT in `slicing.py`'s single threshold block: this is a fact about PDF
#: glyph layout — what counts as one visual LINE — and belongs to whoever reads the PDF.
#: `slicing.py` imports from here, so hosting it there would invert the dependency; and
#: §5's one-block rule exists for the DETECTION thresholds (heading ratio, kept pages,
#: slice fractions) that drifted from their documentation, none of which live here.
_LINE_TOLERANCE_PT = 1.5


def _lines_for_page(page, index: int) -> list[PdfLine]:
    """Group one page's glyphs into visual lines, top-to-bottom then left-to-right.

    Sorted explicitly rather than trusting the content-stream order: a PDF may emit
    glyphs in any order it likes, and a detector fed lines in stream order would
    produce different sections for two files that render identically.
    """
    chars = [c for c in (page.chars or []) if (c.get("text") or "") != ""]
    if not chars:
        return []
    chars.sort(key=lambda c: (round(float(c.get("top") or 0.0), 2), float(c.get("x0") or 0.0)))
    out: list[PdfLine] = []
    group: list[dict] = []
    group_top = float(chars[0].get("top") or 0.0)
    for char in chars:
        top = float(char.get("top") or 0.0)
        if group and abs(top - group_top) > _LINE_TOLERANCE_PT:
            line = _line_from(group, index)
            if line is not None:
                out.append(line)
            group = []
            group_top = top
        group.append(char)
    line = _line_from(group, index)
    if line is not None:
        out.append(line)
    return out


def _line_from(chars: list[dict], index: int) -> PdfLine | None:
    if not chars:
        return None
    text = "".join(str(c.get("text") or "") for c in chars).strip()
    if not text:
        return None
    size = max(float(c.get("size") or 0.0) for c in chars)
    return PdfLine(page=index, text=text, size=round(size, 2), char_count=len(chars))


def _outline_titles(pdf) -> tuple[str, ...]:
    """The PDF's bookmark titles in document order; empty when it has none."""
    try:
        entries = pdf.doc.get_outlines()
    except Exception:
        return ()  # pdfminer raises PDFNoOutlines for the (common) no-bookmarks case
    titles: list[str] = []
    for entry in entries:
        try:
            title = entry[1]
        except (IndexError, TypeError):
            continue
        if isinstance(title, bytes):
            title = title.decode("utf-8", "replace")
        title = " ".join(str(title or "").split())
        if title:
            titles.append(title)
    return tuple(titles)
