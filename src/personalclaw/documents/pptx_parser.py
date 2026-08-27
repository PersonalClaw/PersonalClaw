""".pptx bytes → :class:`DeckModel` + a :class:`LossReport`.

The READ half of the deck round trip. Without it a deck could be written but never loaded,
so the editing surface had nothing to edit and ``GET …/model`` answered 415 for every
.pptx — and a bullet's indent depth could not survive a lap, because the writer pinned
``level = 0`` and nothing ever read one back.

**Depth comes from the file's own attribute**, ``a:pPr/@lvl``, which python-pptx exposes as
``paragraph.level``. Never from the text (leading dashes, indent characters): an outline
whose depth was guessed from prose would flatten the one deck that spells its sub-points
without punctuation, and deepen the one that uses a dash as a word.

**Every parse lives beside a loss report, same contract as the docx and xlsx parsers.** A
real deck holds a great deal this model does not — pictures, tables, charts, free text
boxes, per-character formatting, extra placeholders — and the editor re-renders from the
model, so anything not represented here is gone the moment a save lands. It is reported,
never dropped silently, and located in the terms the EDITOR shows ("slide 3 · bullet 2",
"cover slide") rather than the file's shape count, so a person can find what it names.
"""

from __future__ import annotations

import io
from typing import Any

from personalclaw.documents.docx_parser import LossReport
from personalclaw.documents.model import DECK_LAYOUTS, Bullet, DeckModel, ShapeBox, Slide
from personalclaw.documents.pptx_shapes import body_placeholder, title_index

#: The layout a deck cover uses. A first slide of exactly this shape becomes
#: ``DeckModel.title`` — see :func:`_is_cover`.
_COVER_LAYOUT = "Title Slide"

#: How much of a lost shape's text a report quotes. Enough to recognise it by; a report
#: that pasted a whole slide would bury the other losses.
_SNIP = 80


def parse_pptx(data: bytes) -> tuple[DeckModel, LossReport]:
    """Parse *data* into a :class:`DeckModel` plus everything that did not fit."""
    from pptx import Presentation

    report = LossReport()
    prs = Presentation(io.BytesIO(data))
    slides = list(prs.slides)
    title = ""
    if slides and _is_cover(slides[0]):
        # The writer emits a deck title AS a cover slide, so reading one back into the
        # title field is what closes that loop: a title edited in the editor re-renders
        # into the same cover rather than adding a second one on every save.
        title = _title_text(slides[0])
        _slide(slides.pop(0), "cover slide", report)
    return (
        DeckModel(
            title=title,
            slides=[
                _slide(slide, f"slide {index}", report) for index, slide in enumerate(slides, 1)
            ],
            width_in=_inches(prs.slide_width),
            height_in=_inches(prs.slide_height),
        ),
        report,
    )


def _is_cover(slide: Any) -> bool:
    """Whether *slide* is a deck cover the title field can hold with nothing left over.

    Deliberately narrow. Folding a slide into a field REMOVES it from the slide list, so
    every test here is "would anything on it have nowhere to go?" — a cover carrying
    bullets, speaker notes or a moved title stays a slide, and is edited as one.
    """
    if slide.slide_layout.name != _COVER_LAYOUT or not _title_text(slide):
        return False
    if _notes_text(slide).strip():
        return False
    if any(text.strip() for text, _level in _paragraphs(slide)):
        return False
    return not _box(slide, slide.shapes.title).placed


def _slide(slide: Any, where: str, report: LossReport) -> Slide:
    layout = str(slide.slide_layout.name or "")
    if layout not in DECK_LAYOUTS:
        report.add(
            "slide_layout",
            f"layout {layout!r} is not one this editor can re-create, so a save rebuilds "
            "the slide from its content on a standard layout",
            location=where,
        )
        layout = ""
    bullets = _bullets(slide, where, report)
    _report_shapes(slide, where, report)
    _report_background(slide, where, report)
    return Slide(
        title=_title_text(slide),
        bullets=bullets,
        notes=_notes_text(slide),
        layout=layout,
        title_box=_box(slide, slide.shapes.title),
        body_box=_box(slide, body_placeholder(slide)),
    )


def _paragraphs(slide: Any) -> list[tuple[str, int]]:
    """The body placeholder's paragraphs as (text, depth), before any reporting.

    Split out from :func:`_bullets` so :func:`_is_cover` can ask "is there body content?"
    without a report existing yet — and so the depth is read in exactly one place.
    """
    body = body_placeholder(slide)
    if body is None:
        return []
    return [(str(para.text or ""), int(para.level or 0)) for para in body.text_frame.paragraphs]


def _bullets(slide: Any, where: str, report: LossReport) -> list[Bullet]:
    """The slide's outline, with the character formatting it could not hold reported.

    A TRAILING empty paragraph is dropped rather than carried: PowerPoint leaves one
    behind whenever a placeholder is touched and emptied, and re-rendering it would grow
    a stray bullet glyph on every save. Interior blank lines are kept — those are spacing
    somebody typed on purpose.
    """
    lines = _paragraphs(slide)
    while lines and not lines[-1][0].strip():
        lines.pop()
    body = body_placeholder(slide)
    paragraphs = list(body.text_frame.paragraphs) if body is not None else []
    for index, para in enumerate(paragraphs[: len(lines)], 1):
        styles = _run_styles(para)
        if styles:
            report.add(
                "bullet_run_style",
                f"{'; '.join(styles)} — the model holds a bullet's text and its depth, "
                "not per-character formatting",
                location=f"{where} · bullet {index}",
            )
    return [Bullet(text=text, level=level) for text, level in lines]


def _run_styles(para: Any) -> list[str]:
    """Character formatting inside one bullet, named once each.

    ``None`` means "inherited from the layout", which is the normal case and not a loss —
    so an EXPLICIT off (``bold = False``, a deliberate un-bolding) is reported too. It is
    a decision the file records and the model cannot.
    """
    found: list[str] = []

    def note(text: str) -> None:
        if text not in found:
            found.append(text)

    for run in para.runs:
        font = run.font
        for label, value in (
            ("bold", font.bold),
            ("italic", font.italic),
            ("underline", font.underline),
        ):
            if value is not None and value is not False:
                note(label)
            elif value is False:
                note(f"{label} explicitly turned off")
        if font.size is not None:
            note(f"an explicit size of {font.size.pt:g}pt")
        if font.name:
            note(f"the font {font.name!r}")
        if font.color is not None and font.color.type is not None:
            note("an explicit text colour")
        address = run.hyperlink.address if run.hyperlink is not None else None
        if address:
            note(f"a link to {address}")
    return found


def _report_shapes(slide: Any, where: str, report: LossReport) -> None:
    """Everything on the slide that is not the title, the body, or empty.

    An EMPTY placeholder is skipped: a layout ships them (a Title Slide has a subtitle
    box), and a report that named every unused box would be noise a user learns to
    ignore — which is how the item that matters gets missed.
    """
    carried = {title_index(slide)}
    body = body_placeholder(slide)
    if body is not None:
        carried.add(int(body.placeholder_format.idx))
    for shape in slide.shapes:
        text = str(shape.text_frame.text or "").strip() if shape.has_text_frame else ""
        if shape.is_placeholder:
            if int(shape.placeholder_format.idx) in carried or not text:
                continue
            report.add(
                "slide_placeholder",
                f"the “{shape.name}” placeholder holds text the deck model has no field "
                f"for: {_snip(text)}",
                location=where,
            )
            continue
        detail = f"{shape.name} ({shape.shape_type}) is not carried by the model"
        if text:
            detail = f"{detail}; its text: {_snip(text)}"
        report.add("slide_shape", detail, location=where)


def _report_background(slide: Any, where: str, report: LossReport) -> None:
    from pptx.enum.dml import MSO_FILL

    fill = slide.background.fill
    # BACKGROUND is python-pptx's reading for "nothing was set here, inherit it" — NOT a
    # fill somebody chose. Reporting it would put an item on every slide of every deck.
    if fill.type is not None and fill.type != MSO_FILL.BACKGROUND:
        report.add(
            "slide_feature",
            f"this slide sets its own background ({fill.type}); a save re-renders it with "
            "the template's background",
            location=where,
        )


def _title_text(slide: Any) -> str:
    title = slide.shapes.title
    return str(title.text or "") if title is not None else ""


def _notes_text(slide: Any) -> str:
    # `slide.notes_slide` CREATES a notes slide on first access, so the `has_` check is
    # load-bearing: reading notes must not modify the presentation being read.
    if not slide.has_notes_slide:
        return ""
    return str(slide.notes_slide.notes_text_frame.text or "")


def _box(slide: Any, shape: Any) -> ShapeBox:
    """A shape's geometry, but ONLY when it overrides its layout's.

    A placeholder normally inherits its position, and python-pptx reports the inherited
    value as if the shape declared it. Storing that would pin every shape of every deck at
    whatever the template said the day it was parsed — so the layout's own box is the
    baseline, and an equal reading means "inherited", which the writer leaves alone.
    """
    if shape is None or not shape.is_placeholder:
        return ShapeBox()
    own = (shape.left, shape.top, shape.width, shape.height)
    if any(value is None for value in own):
        return ShapeBox()
    if own == _layout_box(slide, int(shape.placeholder_format.idx)):
        return ShapeBox()
    return ShapeBox(*(_inches(value) for value in own))


def _layout_box(slide: Any, index: int) -> tuple[Any, ...] | None:
    for shape in slide.slide_layout.placeholders:
        if int(shape.placeholder_format.idx) == index:
            return (shape.left, shape.top, shape.width, shape.height)
    return None


def _inches(value: Any) -> float:
    """EMU → inches, rounded to ten-thousandths.

    Rounded because the editor SHOWS this number and "1.4999999999 in" is not a position
    anybody typed; the discarded precision is 45 EMU, about one twenty-thousandth of an
    inch, which no renderer or reader can express.
    """
    return round(float(value.inches), 4)


def _snip(text: str) -> str:
    flat = " ".join(text.split())
    return repr(flat if len(flat) <= _SNIP else flat[: _SNIP - 1] + "…")
