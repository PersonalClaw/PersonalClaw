"""DeckModel → .pptx bytes (python-pptx).

Uses the template's real TITLE + BODY placeholders rather than free-floating text boxes.
That matters for more than tidiness: the pptx reader identifies a slide's title via
``slide.shapes.title``, so a deck built from text boxes would round-trip with every title
lost — and so would PowerPoint's own outline view.

**A bullet's depth is written, not assumed.** This writer used to pin ``level = 0`` on
every appended paragraph, so a nested outline came out of the model flat: the depth was
not lost in the file, it was never written. ``documents/pptx_parser.py`` reads the same
attribute back, which is what makes depth a round trip rather than a one-way flatten.
"""

from __future__ import annotations

import io

from personalclaw.documents.model import DeckModel, ShapeBox, Slide
from personalclaw.documents.pptx_shapes import body_placeholder
from personalclaw.documents.registry import register_writer

#: Layout indexes in python-pptx's default template, used when a slide names no layout.
#: 0 = title slide (title + subtitle), 1 = title and content, 5 = title only.
_LAYOUT_TITLE = 0
_LAYOUT_TITLE_CONTENT = 1
_LAYOUT_TITLE_ONLY = 5


def render_pptx(model: object) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    if not isinstance(model, DeckModel):
        raise TypeError("pptx writer expects a DeckModel")
    prs = Presentation()
    # Slide size first: it is a presentation-level property, so setting it after the
    # slides exist would resize a deck whose shape geometry was placed against the old
    # size. 0.0 leaves the template's own size, per the model's convention.
    if model.width_in > 0:
        prs.slide_width = Inches(model.width_in)
    if model.height_in > 0:
        prs.slide_height = Inches(model.height_in)
    # A deck title becomes a real title slide, so the file opens on something meaningful
    # rather than the first content slide. The parser folds exactly this shape back into
    # `DeckModel.title`, so a title edited in the editor does not add a second cover.
    if model.title:
        layout = prs.slide_layouts[_LAYOUT_TITLE]
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title is not None:
            slide.shapes.title.text = model.title
    for entry in model.slides:
        _add_slide(prs, entry)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _layout_for(prs, entry: Slide):
    """The slide's declared layout, or the one its content asks for.

    Resolved from the REAL template by name rather than from `DECK_LAYOUTS` by index: a
    template whose layouts were renamed or reordered would otherwise silently produce a
    different slide than the one the model names.
    """
    if entry.layout:
        for layout in prs.slide_layouts:
            if layout.name == entry.layout:
                return layout
    # A slide with no body uses the title-only layout: an empty body placeholder renders
    # as a visible "Click to add text" prompt in PowerPoint.
    return prs.slide_layouts[_LAYOUT_TITLE_CONTENT if entry.bullets else _LAYOUT_TITLE_ONLY]


def _add_slide(prs, entry: Slide) -> None:
    slide = prs.slides.add_slide(_layout_for(prs, entry))
    if slide.shapes.title is not None:
        slide.shapes.title.text = entry.title or ""
        _place(slide.shapes.title, entry.title_box)

    if entry.bullets:
        body = body_placeholder(slide)
        if body is not None:
            frame = body.text_frame
            # The first paragraph already exists; reuse it, then append the rest, or the
            # deck opens with a blank first bullet.
            frame.text = entry.bullets[0].text
            frame.paragraphs[0].level = entry.bullets[0].level
            for bullet in entry.bullets[1:]:
                para = frame.add_paragraph()
                para.text = bullet.text
                para.level = bullet.level
            _place(body, entry.body_box)

    # An image block references an artifact; resolving it to bytes is the caller's job
    # (it owns the store). Recording the reference in the notes keeps the fact that an
    # image belonged here rather than dropping it silently.
    notes_text = entry.notes
    if entry.artifact_slug:
        notes_text = (notes_text + f"\n[image: {entry.artifact_slug}]").strip()
    if notes_text:
        # notes_slide is created on first access by python-pptx.
        slide.notes_slide.notes_text_frame.text = notes_text


def _place(shape, box: ShapeBox) -> None:
    """Pin *shape* to *box*, or leave it inheriting the layout's position.

    An unplaced box writes NOTHING: a placeholder with no explicit geometry inherits from
    its layout, and pinning the inherited value would turn "wherever the layout puts it"
    into a frozen position on the first save.
    """
    from pptx.util import Inches

    if not box.placed:
        return
    shape.left = Inches(box.left_in)
    shape.top = Inches(box.top_in)
    shape.width = Inches(box.width_in)
    shape.height = Inches(box.height_in)


register_writer("pptx", render_pptx)
